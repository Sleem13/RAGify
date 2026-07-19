from langchain_community.document_loaders import Docx2txtLoader, TextLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document
import os
import tempfile
import logging
import base64

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {
    ".docx": Docx2txtLoader,
    ".txt":  TextLoader,
}

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg"}
PPT_EXTENSIONS = {".pptx"}
PDF_EXTENSIONS = {".pdf"}

class DocumentProcessor:
    """
    Handles loading, chunking, and preparing documents for vector storage.
    Supports PDF (with embedded images via PyMuPDF + Gemini Vision), 
    DOCX, TXT, PPTX, and standalone Images.
    """

    def __init__(self):
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1200,
            chunk_overlap=250,
            length_function=len,
        )

    async def process_file(self, file_bytes: bytes, filename: str):
        ext = os.path.splitext(filename)[1].lower()

        if not self.is_supported(filename):
            raise ValueError(
                f"Unsupported file type '{ext}'. Supported: PDF, DOCX, TXT, PPTX, PNG, JPG, JPEG."
            )

        documents = []

        # ── 1. Image Processing (via Gemini Vision) ───────────────────────────
        if ext in IMAGE_EXTENSIONS:
            extracted_text = await self._extract_image_text(file_bytes, ext)
            documents = [Document(page_content=extracted_text, metadata={"source": filename})]

        # ── 2. PowerPoint Processing ──────────────────────────────────────────
        elif ext in PPT_EXTENSIONS:
            extracted_text = self._extract_pptx_text(file_bytes)
            documents = [Document(page_content=extracted_text, metadata={"source": filename})]

        # ── 3. Advanced PDF Processing (PyMuPDF + Vision for embedded images) ─
        elif ext in PDF_EXTENSIONS:
            extracted_text = await self._extract_pdf_text_and_images(file_bytes, filename)
            documents = [Document(page_content=extracted_text, metadata={"source": filename})]

        # ── 4. Standard Document Processing (DOCX, TXT) ───────────────────────
        else:
            loader_class = SUPPORTED_EXTENSIONS[ext]
            temp_file_path = None
            with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
                tmp.write(file_bytes)
                temp_file_path = tmp.name

            try:
                if ext == ".txt":
                    loader = loader_class(temp_file_path, encoding="utf-8")
                else:
                    loader = loader_class(temp_file_path)
                documents = loader.load()
                for doc in documents:
                    doc.metadata["source"] = filename
            finally:
                if temp_file_path and os.path.exists(temp_file_path):
                    os.remove(temp_file_path)

        # ── Chunking ──────────────────────────────────────────────────────────
        chunks = self.text_splitter.split_documents(documents)
        logger.info(
            f"Processed '{filename}': {len(documents)} page(s) → {len(chunks)} chunks."
        )
        return chunks

    async def _extract_pdf_text_and_images(self, file_bytes: bytes, filename: str) -> str:
        """
        Uses PyMuPDF (fitz) to extract text and embedded images from a PDF.
        Sends images to Gemini Vision to describe charts/tables.
        """
        try:
            import fitz
        except ImportError:
            raise ValueError("PyMuPDF (fitz) is not installed. Run: pip install pymupdf")

        pdf_document = fitz.open(stream=file_bytes, filetype="pdf")
        full_text = []

        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            page_text = page.get_text("text").strip()
            
            # Append page text
            full_text.append(f"--- Page {page_num + 1} ---")
            if page_text:
                full_text.append(page_text)
            
            # Extract images from this page
            image_list = page.get_images(full=True)
            for img_index, img in enumerate(image_list):
                xref = img[0]
                base_image = pdf_document.extract_image(xref)
                image_bytes = base_image["image"]
                img_ext = base_image["ext"]
                
                # Only process reasonable sized images (skip tiny icons)
                if len(image_bytes) > 5000:  
                    try:
                        logger.info(f"Extracting vision data from image on page {page_num + 1}")
                        vision_text = await self._extract_image_text(image_bytes, f".{img_ext}")
                        full_text.append(f"\n[Image/Chart detected on page {page_num + 1}]:\n{vision_text}\n")
                    except Exception as e:
                        logger.warning(f"Failed to process image on page {page_num + 1}: {e}")

        pdf_document.close()
        return "\n".join(full_text)

    async def _extract_image_text(self, file_bytes: bytes, ext: str) -> str:
        """Extract text from an image using Gemini Vision."""
        gemini_key = os.getenv("GEMINI_API_KEY", "").strip()
        if not gemini_key or gemini_key == "your_gemini_key_here":
            logger.warning("No GEMINI_API_KEY found, skipping image vision processing.")
            return "Image analysis skipped: No Gemini API Key configured."

        from langchain_google_genai import ChatGoogleGenerativeAI
        from langchain_core.messages import HumanMessage
        
        b64_img = base64.b64encode(file_bytes).decode('utf-8')
        mime_type = f"image/{ext.strip('.')}"
        if mime_type == "image/jpg": 
            mime_type = "image/jpeg"

        llm = ChatGoogleGenerativeAI(model="gemini-1.5-flash", google_api_key=gemini_key)
        msg = HumanMessage(content=[
            {
                "type": "text", 
                "text": "Analyze this image in detail. Extract any readable text. If it is a chart, graph, or table, describe the data points, trends, and key takeaways clearly so that it can be searched."
            },
            {
                "type": "image_url", 
                "image_url": {"url": f"data:{mime_type};base64,{b64_img}"}
            }
        ])
        
        response = await llm.ainvoke([msg])
        return response.content

    def _extract_pptx_text(self, file_bytes: bytes) -> str:
        """Extract text from a PowerPoint presentation."""
        try:
            import io
            from pptx import Presentation
        except ImportError:
            raise ValueError("python-pptx is not installed. Run: pip install python-pptx")

        prs = Presentation(io.BytesIO(file_bytes))
        text_runs = []
        for slide_idx, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if slide_text:
                text_runs.append(f"--- Slide {slide_idx + 1} ---\n" + "\n".join(slide_text))
        
        return "\n\n".join(text_runs)

    @staticmethod
    def is_supported(filename: str) -> bool:
        ext = os.path.splitext(filename)[1].lower()
        return ext in SUPPORTED_EXTENSIONS or ext in IMAGE_EXTENSIONS or ext in PPT_EXTENSIONS or ext in PDF_EXTENSIONS

# Singleton
document_processor = DocumentProcessor()
