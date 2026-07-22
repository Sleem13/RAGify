from langchain_community.document_loaders import Docx2txtLoader, TextLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document
import asyncio
import os
from pathlib import Path
import tempfile
import logging
import base64

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {
    ".docx": Docx2txtLoader,
    ".txt":  TextLoader,
}

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg"}
PPT_EXTENSIONS = {".pptx"}
PDF_EXTENSIONS = {".pdf"}

class DocumentProcessor:
    """
    Handles loading, chunking, and preparing documents for vector storage.
    Supports PDF (with embedded images via PyMuPDF + Gemini Vision), 
    DOCX, TXT, PPTX, and standalone Images.
    """

    def __init__(self):
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1200,
            chunk_overlap=250,
            length_function=len,
        )
        self.max_pdf_images = self._env_int("RAGIFY_MAX_PDF_IMAGES", 4)
        self.pdf_image_timeout = self._env_int("RAGIFY_PDF_IMAGE_TIMEOUT", 20)
        self.pdf_image_concurrency = self._env_int("RAGIFY_PDF_IMAGE_CONCURRENCY", 2)

    @staticmethod
    def _env_int(name: str, default: int) -> int:
        try:
            return max(0, int(os.getenv(name, str(default))))
        except ValueError:
            logger.warning("Invalid %s value; using %s.", name, default)
            return default

    async def process_file(self, file_bytes: bytes, filename: str):
        ext = Path(filename).suffix.lower()

        if not self.is_supported(filename):
            raise ValueError(
                f"Unsupported file type '{ext}'. Supported: PDF, DOCX, TXT, PPTX, PNG, JPG, JPEG."
            )

        documents = []

        # ── 1. Image Processing (via Gemini Vision) ───────────────────────────
        if ext in IMAGE_EXTENSIONS:
            extracted_text = await self._extract_image_text(file_bytes, ext)
            documents = [Document(page_content=extracted_text, metadata={"source": filename})]

        # ── 2. PowerPoint Processing ──────────────────────────────────────────
        elif ext in PPT_EXTENSIONS:
            extracted_text = self._extract_pptx_text(file_bytes)
            documents = [Document(page_content=extracted_text, metadata={"source": filename})]

        # ── 3. Advanced PDF Processing (PyMuPDF + Vision for embedded images) ─
        elif ext in PDF_EXTENSIONS:
            documents = await self._extract_pdf_documents(file_bytes, filename)

        # ── 4. Standard Document Processing (DOCX, TXT) ───────────────────────
        else:
            loader_class = SUPPORTED_EXTENSIONS[ext]
            temp_file_path: Path | None = None
            with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
                tmp.write(file_bytes)
                temp_file_path = Path(tmp.name)

            try:
                if ext == ".txt":
                    loader = loader_class(str(temp_file_path), encoding="utf-8")
                else:
                    loader = loader_class(str(temp_file_path))
                documents = loader.load()
                for doc in documents:
                    doc.metadata["source"] = filename
            finally:
                if temp_file_path:
                    temp_file_path.unlink(missing_ok=True)

        # ── Chunking ──────────────────────────────────────────────────────────
        documents = [doc for doc in documents if doc.page_content.strip()]
        chunks = self.text_splitter.split_documents(documents)
        for chunk_index, chunk in enumerate(chunks):
            chunk.metadata.update(
                {"source": filename, "chunk_index": chunk_index, "is_current": True}
            )
        logger.info(
            f"Processed '{filename}': {len(documents)} page(s) → {len(chunks)} chunks."
        )
        return chunks

    async def _extract_pdf_documents(
        self, file_bytes: bytes, filename: str
    ) -> list[Document]:
        """
        Uses PyMuPDF (fitz) to extract text and embedded images from a PDF.
        Sends images to Gemini Vision to describe charts/tables.
        """
        try:
            import fitz
        except ImportError:
            raise ValueError("PyMuPDF (fitz) is not installed. Run: pip install pymupdf")

        def extract_pages() -> tuple[list[tuple[int, str]], list[tuple[int, int, bytes, str]]]:
            pages: list[tuple[int, str]] = []
            candidates: list[tuple[int, int, bytes, str]] = []
            seen_xrefs: set[int] = set()
            with fitz.open(stream=file_bytes, filetype="pdf") as pdf_document:
                for page_num in range(len(pdf_document)):
                    page = pdf_document.load_page(page_num)
                    page_text = page.get_text("text").strip()
                    pages.append((page_num + 1, page_text))
                    for image in page.get_images(full=True):
                        xref = image[0]
                        if xref in seen_xrefs:
                            continue
                        seen_xrefs.add(xref)
                        base_image = pdf_document.extract_image(xref)
                        image_bytes = base_image["image"]
                        if len(image_bytes) > 5000:
                            candidates.append(
                                (page_num + 1, len(page_text), image_bytes, base_image["ext"])
                            )
            return pages, candidates

        pages, image_candidates = await asyncio.to_thread(extract_pages)

        image_candidates.sort(key=lambda item: (item[1] >= 120, -len(item[2])))
        selected_images = image_candidates[: self.max_pdf_images]
        if len(image_candidates) > len(selected_images):
            logger.info(
                "PDF '%s' has %s unique image candidates; analyzing the top %s.",
                filename,
                len(image_candidates),
                len(selected_images),
            )

        semaphore = asyncio.Semaphore(max(1, self.pdf_image_concurrency))

        async def analyze_image(candidate: tuple[int, int, bytes, str]):
            page_number, _, image_bytes, image_extension = candidate
            try:
                async with semaphore:
                    vision_text = await asyncio.wait_for(
                        self._extract_image_text(image_bytes, f".{image_extension}"),
                        timeout=max(1, self.pdf_image_timeout),
                    )
                return page_number, vision_text
            except Exception as exc:
                logger.warning("Image analysis failed on page %s: %s", page_number, exc)
                return page_number, ""

        vision_by_page: dict[int, list[str]] = {}
        if selected_images:
            vision_results = await asyncio.gather(
                *(analyze_image(candidate) for candidate in selected_images)
            )
            for page_number, vision_text in vision_results:
                if vision_text:
                    vision_by_page.setdefault(page_number, []).append(vision_text)

        documents: list[Document] = []
        for page_number, page_text in pages:
            parts = [page_text]
            parts.extend(
                f"[Image or chart on this page]\n{vision_text}"
                for vision_text in vision_by_page.get(page_number, [])
            )
            content = "\n\n".join(part for part in parts if part.strip()).strip()
            if content:
                documents.append(
                    Document(
                        page_content=content,
                        metadata={"source": filename, "page": page_number},
                    )
                )
        return documents

    async def _extract_image_text(self, file_bytes: bytes, ext: str) -> str:
        """Extract text from an image using Gemini Vision."""
        gemini_key = os.getenv("GEMINI_API_KEY", "").strip()
        if not gemini_key or gemini_key == "your_gemini_key_here":
            logger.warning("No GEMINI_API_KEY found, skipping image vision processing.")
            raise ValueError("Image analysis requires a configured GEMINI_API_KEY.")

        from langchain_google_genai import ChatGoogleGenerativeAI
        from langchain_core.messages import HumanMessage
        
        b64_img = base64.b64encode(file_bytes).decode('utf-8')
        mime_type = f"image/{ext.strip('.')}"
        if mime_type == "image/jpg": 
            mime_type = "image/jpeg"

        llm = ChatGoogleGenerativeAI(model="gemini-2.5-flash", google_api_key=gemini_key)
        msg = HumanMessage(content=[
            {
                "type": "text", 
                "text": "Analyze this image in detail. Extract any readable text. If it is a chart, graph, or table, describe the data points, trends, and key takeaways clearly so that it can be searched."
            },
            {
                "type": "image_url", 
                "image_url": {"url": f"data:{mime_type};base64,{b64_img}"}
            }
        ])
        
        response = await llm.ainvoke([msg])
        return response.content

    def _extract_pptx_text(self, file_bytes: bytes) -> str:
        """Extract text from a PowerPoint presentation."""
        try:
            import io
            from pptx import Presentation
        except ImportError:
            raise ValueError("python-pptx is not installed. Run: pip install python-pptx")

        prs = Presentation(io.BytesIO(file_bytes))
        text_runs = []
        for slide_idx, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if slide_text:
                text_runs.append(f"--- Slide {slide_idx + 1} ---\n" + "\n".join(slide_text))
        
        return "\n\n".join(text_runs)

    @staticmethod
    def is_supported(filename: str) -> bool:
        ext = Path(filename).suffix.lower()
        return ext in SUPPORTED_EXTENSIONS or ext in IMAGE_EXTENSIONS or ext in PPT_EXTENSIONS or ext in PDF_EXTENSIONS

    async def process_excel_for_rag(self, file_bytes: bytes, filename: str):
        """
        Convert Excel/CSV data into text chunks for FAISS indexing.
        Enables Q&A over tabular data (e.g., 'What is the highest sale?')
        """
        import pandas as pd
        import tempfile

        ext = Path(filename).suffix.lower()
        temp_file_path: Path | None = None

        with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
            tmp.write(file_bytes)
            temp_file_path = Path(tmp.name)

        try:
            if ext == ".csv":
                df = pd.read_csv(temp_file_path, encoding="utf-8", on_bad_lines="skip")
            elif ext == ".json":
                try:
                    df = pd.read_json(temp_file_path)
                except ValueError:
                    df = pd.read_json(temp_file_path, lines=True)
            elif ext == ".xls":
                df = pd.read_excel(temp_file_path, engine="xlrd")
            else:
                df = pd.read_excel(temp_file_path, engine="openpyxl")

            df = df.where(pd.notnull(df), None)
            documents = []

            # Create a text summary of the entire dataset
            columns = list(df.columns)
            numeric_cols = df.select_dtypes(include=["number"]).columns.tolist()

            overview = (
                f"File: {filename}\n"
                f"Total rows: {len(df)}\n"
                f"Columns: {', '.join(columns)}\n"
            )
            if numeric_cols:
                stats = df[numeric_cols].describe().to_string()
                overview += f"\nStatistical Summary:\n{stats}\n"

            documents.append(Document(
                page_content=overview,
                metadata={"source": filename, "chunk_type": "overview"}
            ))

            # Create chunks of rows (every 20 rows = 1 chunk)
            chunk_size = 20
            for start in range(0, len(df), chunk_size):
                chunk_df = df.iloc[start:start + chunk_size]
                chunk_text = f"Data from '{filename}' (rows {start+1} to {start+len(chunk_df)}):\n"
                chunk_text += chunk_df.to_string(index=False)
                documents.append(Document(
                    page_content=chunk_text,
                    metadata={"source": filename, "chunk_type": "data_rows", "row_start": start}
                ))

            # Create summary stats per column
            for col in numeric_cols:
                col_stats = (
                    f"Column '{col}' statistics from '{filename}':\n"
                    f"  Min: {df[col].min():.2f}\n"
                    f"  Max: {df[col].max():.2f}\n"
                    f"  Average: {df[col].mean():.2f}\n"
                    f"  Sum: {df[col].sum():.2f}\n"
                    f"  Median: {df[col].median():.2f}\n"
                )
                documents.append(Document(
                    page_content=col_stats,
                    metadata={"source": filename, "chunk_type": "column_stats", "column": col}
                ))

            logger.info(f"Excel '{filename}' converted to {len(documents)} RAG chunks.")
            return documents

        except Exception as exc:
            logger.error(f"Failed to convert Excel to RAG: {exc}")
            raise ValueError(f"Failed to convert spreadsheet for retrieval: {exc}") from exc
        finally:
            if temp_file_path:
                temp_file_path.unlink(missing_ok=True)

# Singleton
document_processor = DocumentProcessor()
